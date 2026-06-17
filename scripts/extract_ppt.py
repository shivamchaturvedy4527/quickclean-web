#!/usr/bin/env python3
"""Extract all text and images from PPTX slides."""
import json
import os
import re
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPT_PATH = r"C:\Users\shiva\OneDrive\Desktop\UNITED MACHINES PROFILE ppt.pptx"
OUT_DIR = Path(__file__).resolve().parent.parent / "public" / "images" / "ppt"
SUMMARY_PATH = Path(__file__).resolve().parent.parent / "data" / "ppt-extraction.json"


def slugify(text, max_len=40):
    text = re.sub(r"[^\w\s-]", "", text.lower())
    text = re.sub(r"[\s_]+", "-", text).strip("-")
    return (text or "image")[:max_len]


def get_shape_text(shape):
    texts = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            line = "".join(run.text for run in para.runs).strip()
            if line:
                level = para.level
                prefix = "  " * level + ("• " if level > 0 else "")
                texts.append(prefix + line)
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        table = shape.table
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                texts.append(" | ".join(cells))
    return texts


def extract_images_from_slide(slide, slide_num, img_counter):
    saved = []
    for shape_idx, shape in enumerate(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img = shape.image
            ext = img.ext or "png"
            blob = img.blob
            name_hint = ""
            if shape.has_text_frame:
                name_hint = shape.text_frame.text.strip()
            fname = f"slide{slide_num:02d}_img{shape_idx:02d}"
            if name_hint:
                fname += f"_{slugify(name_hint)}"
            fname += f".{ext}"
            out_path = OUT_DIR / fname
            with open(out_path, "wb") as f:
                f.write(blob)
            saved.append({
                "filename": fname,
                "path": f"/images/ppt/{fname}",
                "size": len(blob),
                "ext": ext,
            })
            img_counter[0] += 1
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for gshape in shape.shapes:
                if gshape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = gshape.image
                    ext = img.ext or "png"
                    blob = img.blob
                    fname = f"slide{slide_num:02d}_group{shape_idx:02d}_img{img_counter[0]:03d}.{ext}"
                    out_path = OUT_DIR / fname
                    with open(out_path, "wb") as f:
                        f.write(blob)
                    saved.append({
                        "filename": fname,
                        "path": f"/images/ppt/{fname}",
                        "size": len(blob),
                        "ext": ext,
                    })
                    img_counter[0] += 1
    return saved


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation(PPT_PATH)
    slides_data = []
    img_counter = [0]

    print(f"Total slides: {len(prs.slides)}")

    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            texts.extend(get_shape_text(shape))

        images = extract_images_from_slide(slide, i, img_counter)

        slide_info = {
            "slide": i,
            "text": texts,
            "text_joined": "\n".join(texts),
            "images": images,
            "image_count": len(images),
        }
        slides_data.append(slide_info)
        print(f"\n=== SLIDE {i} ({len(images)} images) ===")
        print(slide_info["text_joined"] or "(no text)")

    # Also extract all media from zip for completeness
    zip_images = []
    with zipfile.ZipFile(PPT_PATH, "r") as zf:
        for name in zf.namelist():
            if name.startswith("ppt/media/"):
                data = zf.read(name)
                base = os.path.basename(name)
                out_path = OUT_DIR / f"media_{base}"
                if not out_path.exists():
                    with open(out_path, "wb") as f:
                        f.write(data)
                zip_images.append({
                    "zip_path": name,
                    "saved_as": f"/images/ppt/media_{base}",
                    "size": len(data),
                })

    result = {
        "source": PPT_PATH,
        "total_slides": len(slides_data),
        "total_slide_images": img_counter[0],
        "total_zip_media": len(zip_images),
        "slides": slides_data,
        "zip_media": zip_images,
    }

    with open(SUMMARY_PATH, "w", encoding="utf-8") as f:
        json.dump(result, f, indent=2, ensure_ascii=False)

    print(f"\n\nExtraction complete:")
    print(f"  Slides: {len(slides_data)}")
    print(f"  Slide images: {img_counter[0]}")
    print(f"  Zip media files: {len(zip_images)}")
    print(f"  Summary: {SUMMARY_PATH}")


if __name__ == "__main__":
    main()
